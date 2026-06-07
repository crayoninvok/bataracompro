"""
Generate PT Batara Dharma Persada — Projects Portfolio PPTX
Run: python scripts/generate-proyek-pptx.py
Output: presentations/PT-BDP-Projects-Portfolio.pptx
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
PUBLIC = ROOT / "public"
OUT_DIR = ROOT / "presentations"
OUT_FILE = OUT_DIR / "PT-BDP-Projects-Portfolio.pptx"

# Brand palette (matches website)
TEAL = RGBColor(0x1F, 0xBF, 0xB8)
ORANGE = RGBColor(0xE8, 0x5C, 0x23)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF5)
MID_GRAY = RGBColor(0xE5, 0xE5, 0xE5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def img(path: str) -> Path:
    return PUBLIC / path.lstrip("/")


def set_run(run, *, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_header_bar(slide, title: str, subtitle: str | None = None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), DARK)
    add_rect(slide, Inches(0), Inches(1.05), SLIDE_W, Inches(0.06), TEAL)

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(10), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.62), Inches(11), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run(sr, size=12, color=TEAL)


def add_footer(slide, text: str = "PT Batara Dharma Persada — Confidential"):
    add_rect(slide, Inches(0), Inches(7.12), SLIDE_W, Inches(0.38), LIGHT_GRAY)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.16), Inches(12), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, size=9, color=GRAY)


def add_bullets(slide, items: list[str], left, top, width, height, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = item
        set_run(r, size=size, color=DARK)
        p.bullet = True


def add_stat_cards(slide, stats: list[tuple[str, str]], left, top):
    card_w = Inches(2.75)
    card_h = Inches(1.15)
    gap = Inches(0.2)
    for i, (value, label) in enumerate(stats):
        x = left + i * (card_w + gap)
        card = add_rect(slide, x, top, card_w, card_h, LIGHT_GRAY)
        card.line.color.rgb = MID_GRAY

        vbox = slide.shapes.add_textbox(x + Inches(0.15), top + Inches(0.12), card_w - Inches(0.3), Inches(0.45))
        vtf = vbox.text_frame
        vtf.clear()
        vp = vtf.paragraphs[0]
        vr = vp.add_run()
        vr.text = value
        set_run(vr, size=22, bold=True, color=ORANGE)

        lbox = slide.shapes.add_textbox(x + Inches(0.15), top + Inches(0.58), card_w - Inches(0.3), Inches(0.4))
        ltf = lbox.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label.upper()
        set_run(lr, size=9, bold=True, color=GRAY)


def add_picture_safe(slide, path: Path, left, top, width, height):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        placeholder = add_rect(slide, left, top, width, height, LIGHT_GRAY)
        placeholder.line.color.rgb = MID_GRAY


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, Inches(0), Inches(3.35), SLIDE_W, Inches(0.08), TEAL)
    add_rect(slide, Inches(0), Inches(3.43), SLIDE_W, Inches(0.04), ORANGE)

    logo_path = img("nobgbtrlogo.png")
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.75), Inches(0.65), width=Inches(1.1))

    title = slide.shapes.add_textbox(Inches(0.75), Inches(3.75), Inches(11.5), Inches(1.2))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Projects Portfolio"
    set_run(r, size=40, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(4.75), Inches(11), Inches(0.8))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Coal Hauling Contractor — PT Batara Dharma Persada"
    set_run(sr, size=18, color=TEAL)

    tag = slide.shapes.add_textbox(Inches(0.75), Inches(5.55), Inches(11), Inches(0.5))
    ttf = tag.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tr = tp.add_run()
    tr.text = "Trusted partner for PT Indonesia Pratama · Bayan Resources Group"
    set_run(tr, size=13, color=GRAY)

    date_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.55), Inches(5), Inches(0.35))
    dtf = date_box.text_frame
    dtf.clear()
    dp = dtf.paragraphs[0]
    dr = dp.add_run()
    dr.text = date.today().strftime("%B %Y")
    set_run(dr, size=11, color=GRAY)

    if img("senyiur/sny2.jpeg").exists():
        add_picture_safe(slide, img("senyiur/sny2.jpeg"), Inches(8.2), Inches(0.55), Inches(4.6), Inches(2.55))


def slide_section(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, Inches(0.75), Inches(3.1), Inches(1.2), Inches(0.07), ORANGE)

    box = slide.shapes.add_textbox(Inches(0.75), Inches(3.35), Inches(11), Inches(1))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run(r, size=34, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(4.35), Inches(10), Inches(0.6))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = subtitle
    set_run(sr, size=15, color=TEAL)


def slide_overview(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Company Overview", "Dedicated coal hauling contractor since 2024")
    add_footer(slide)

    bullets = [
        "Specialized coal hauling contractor — pit to stockpile and designated delivery points",
        "Active operations at Senyiur, East Kalimantan for PT Indonesia Pratama",
        "Muara Pahu project in preparation — go-live 2nd week of July 2026",
        "300++ workforce · 30 active hauling units · 5 million ton annual capacity",
        "Volvo FH 16-700 fleet with double-trailer configuration (110 m³ + 125 m³)",
    ]
    add_bullets(slide, bullets, Inches(0.65), Inches(1.45), Inches(6.2), Inches(5.2), size=15)

    add_stat_cards(
        slide,
        [("30", "Active Fleet"), ("5M", "Tons / Year"), ("300++", "Employees"), ("2024", "Since")],
        Inches(7.2),
        Inches(1.55),
    )

    add_picture_safe(slide, img("nobgbtr.png"), Inches(7.35), Inches(3.35), Inches(4.8), Inches(2.8))


def slide_client(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Client Partnership", "PT Indonesia Pratama · PT Bayan Resources Tbk")
    add_footer(slide)

    add_rect(slide, Inches(0.65), Inches(1.55), Inches(5.8), Inches(4.85), LIGHT_GRAY)
    add_rect(slide, Inches(0.65), Inches(1.55), Inches(5.8), Inches(0.55), TEAL)
    client_title = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(5.4), Inches(0.4))
    ctf = client_title.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cr = cp.add_run()
    cr.text = "OUR CLIENT"
    set_run(cr, size=11, bold=True, color=WHITE)

    logo_box = add_rect(slide, Inches(1.5), Inches(2.45), Inches(3.1), Inches(1.6), WHITE)
    logo_box.line.color.rgb = MID_GRAY
    add_picture_safe(slide, img("bayan/ipbayan.png"), Inches(1.75), Inches(2.65), Inches(2.6), Inches(1.2))

    name = slide.shapes.add_textbox(Inches(0.95), Inches(4.35), Inches(5.2), Inches(0.45))
    ntf = name.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.CENTER
    nr = np.add_run()
    nr.text = "PT Indonesia Pratama"
    set_run(nr, size=18, bold=True, color=DARK)

    group = slide.shapes.add_textbox(Inches(0.95), Inches(4.78), Inches(5.2), Inches(0.35))
    gtf = group.text_frame
    gtf.clear()
    gp = gtf.paragraphs[0]
    gp.alignment = PP_ALIGN.CENTER
    gr = gp.add_run()
    gr.text = "Member of PT Bayan Resources Group"
    set_run(gr, size=12, color=GRAY)

    add_bullets(
        slide,
        [
            "Long-term coal hauling contract at Senyiur since 2024",
            "Expanding partnership to Muara Pahu project site in 2026",
            "Aligned with client HSE, production, and reporting standards",
            "Proven track record of scalable fleet deployment",
        ],
        Inches(6.85),
        Inches(1.65),
        Inches(5.8),
        Inches(4.5),
        size=14,
    )


def slide_senyiur(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Senyiur — Active Project", "Coal hauling contract · East Kalimantan")
    add_footer(slide)

    add_stat_cards(
        slide,
        [("30 Units", "Active Fleet"), ("5M Tons", "Annual Capacity"), ("300++", "Workforce"), ("2024", "Since")],
        Inches(0.65),
        Inches(1.35),
    )

    scope = [
        "Coal hauling from pit to stockpile and designated delivery points",
        "Double-trailer fleet with Volvo FH 16-700 tractor heads",
        "On-site dispatch, route management, and production support",
        "In-house workshop maintenance for fleet readiness and uptime",
        "HSE-compliant operations aligned with client and regulatory standards",
    ]
    add_bullets(slide, scope, Inches(0.65), Inches(2.85), Inches(6.3), Inches(3.8), size=13)

    add_picture_safe(slide, img("senyiur/sny1.jpeg"), Inches(7.15), Inches(1.35), Inches(5.55), Inches(3.15))
    cap = slide.shapes.add_textbox(Inches(7.15), Inches(4.55), Inches(5.55), Inches(0.35))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "Senyiur Mining Site — Active Coal Hauling Operations"
    set_run(cr, size=10, bold=True, color=GRAY)


def slide_fleet(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Fleet Specification", "Heavy-duty coal hauling configuration")
    add_footer(slide)

    specs = [
        ("Tractor Head", "Volvo FH 16-700 HP"),
        ("Trailer Configuration", "Double Vessel (Side-Tip)"),
        ("Payload Capacity", "110 m³ + 125 m³"),
        ("Service Type", "Coal Hauling Contractor"),
        ("Operation Model", "Pit to Stockpile / Hopper"),
        ("Maintenance", "In-house Batara Plant Workshop"),
    ]

    y = Inches(1.55)
    for label, value in specs:
        add_rect(slide, Inches(0.75), y, Inches(11.85), Inches(0.72), LIGHT_GRAY if specs.index((label, value)) % 2 == 0 else WHITE)
        lbox = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.16), Inches(4.5), Inches(0.4))
        ltf = lbox.text_frame
        ltf.clear()
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        set_run(lr, size=13, color=GRAY)

        vbox = slide.shapes.add_textbox(Inches(5.8), y + Inches(0.14), Inches(6.5), Inches(0.45))
        vtf = vbox.text_frame
        vtf.clear()
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.RIGHT
        vr = vp.add_run()
        vr.text = value
        set_run(vr, size=14, bold=True, color=DARK)
        y += Inches(0.78)

    add_picture_safe(slide, img("senyiur/sny4.jpeg"), Inches(8.55), Inches(1.45), Inches(4.05), Inches(2.35))


def slide_hauling(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Coal Hauling Expertise", "What we deliver on every project site")
    add_footer(slide)

    services = [
        ("Pit-to-Stockpile Hauling", "Dedicated routes with consistent cycle times and load accuracy"),
        ("Fleet Dispatch & Control", "Shift planning and real-time monitoring to meet production targets"),
        ("Workshop & Maintenance", "Preventive maintenance and rapid breakdown response"),
        ("HSE & Compliance", "Safety-first operations with driver training and regulatory compliance"),
        ("Production Reporting", "Daily tonnage tracking and KPI reporting for transparent management"),
        ("24/7 Operation Support", "Round-the-clock hauling aligned with mine production schedules"),
    ]

    col_w = Inches(5.85)
    for i, (title, desc) in enumerate(services):
        col = i % 2
        row = i // 2
        x = Inches(0.65) + col * (col_w + Inches(0.35))
        y = Inches(1.45) + row * Inches(1.75)

        card = add_rect(slide, x, y, col_w, Inches(1.55), LIGHT_GRAY)
        card.line.color.rgb = MID_GRAY
        add_rect(slide, x, y, Inches(0.08), Inches(1.55), TEAL if i % 2 == 0 else ORANGE)

        tbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.4), Inches(0.35))
        ttf = tbox.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        set_run(tr, size=13, bold=True, color=DARK)

        dbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.5), col_w - Inches(0.4), Inches(0.95))
        dtf = dbox.text_frame
        dtf.word_wrap = True
        dtf.clear()
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = desc
        set_run(dr, size=11, color=GRAY)


def slide_muara_pahu(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Muara Pahu — Upcoming Project", "In preparation · Go-live 2nd week of July 2026")
    add_footer(slide)

    add_stat_cards(
        slide,
        [("23 Units", "Planned Fleet"), ("Jul 2026", "Go-Live"), ("2nd Week", "Target Start"), ("Prep", "Status")],
        Inches(0.65),
        Inches(1.35),
    )

    prep = [
        "Fleet mobilization and equipment staging for 23 hauling units",
        "Manpower recruitment, deployment, and site induction programs",
        "Workshop and support facility setup at Muara Pahu project site",
        "Route planning, dispatch systems, and operational readiness review",
        "HSE briefing and pre-commissioning aligned with PT Indonesia Pratama standards",
    ]
    add_bullets(slide, prep, Inches(0.65), Inches(2.85), Inches(6.2), Inches(3.5), size=13)

    loc = slide.shapes.add_textbox(Inches(0.65), Inches(6.35), Inches(6.5), Inches(0.45))
    ltf = loc.text_frame
    ltf.clear()
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = "Coordinates: 0°01'24.3\"N 116°01'45.2\"E · Muara Pahu, East Kalimantan"
    set_run(lr, size=11, bold=True, color=ORANGE)

    add_picture_safe(slide, img("muarapahu/mhu1.jpeg"), Inches(7.15), Inches(1.35), Inches(5.55), Inches(3.15))
    cap = slide.shapes.add_textbox(Inches(7.15), Inches(4.55), Inches(5.55), Inches(0.35))
    ctf = cap.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "Muara Pahu — Site Preparation Underway"
    set_run(cr, size=10, bold=True, color=GRAY)


def slide_photos_senyiur(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Senyiur Operations", "Field documentation — PT Indonesia Pratama")
    add_footer(slide)

    photos = [
        ("senyiur/sny3.jpeg", "Pit to Stockpile Route"),
        ("senyiur/sny8.jpeg", "Hopper Coal Discharge"),
        ("senyiur/sny6.JPG", "Site Operations"),
        ("senyiur/sny7.JPG", "Hauling Fleet"),
    ]
    positions = [
        (Inches(0.65), Inches(1.4)),
        (Inches(6.95), Inches(1.4)),
        (Inches(0.65), Inches(4.05)),
        (Inches(6.95), Inches(4.05)),
    ]
    for (path, caption), (x, y) in zip(photos, positions):
        add_picture_safe(slide, img(path), x, y, Inches(5.95), Inches(2.35))
        cbox = slide.shapes.add_textbox(x, y + Inches(2.42), Inches(5.95), Inches(0.3))
        ctf = cbox.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = caption
        set_run(cr, size=10, bold=True, color=GRAY)


def slide_photos_muara(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Muara Pahu Preparation", "Progress ahead of July 2026 go-live")
    add_footer(slide)

    photos = [
        ("muarapahu/mhu2.jpeg", "Fleet Mobilization"),
        ("muarapahu/mhu3.jpeg", "Workshop Setup"),
        ("muarapahu/mhu5.jpeg", "Site Development"),
        ("muarapahu/mhu6.jpeg", "Preparation Progress"),
    ]
    positions = [
        (Inches(0.65), Inches(1.4)),
        (Inches(6.95), Inches(1.4)),
        (Inches(0.65), Inches(4.05)),
        (Inches(6.95), Inches(4.05)),
    ]
    for (path, caption), (x, y) in zip(photos, positions):
        add_picture_safe(slide, img(path), x, y, Inches(5.95), Inches(2.35))
        cbox = slide.shapes.add_textbox(x, y + Inches(2.42), Inches(5.95), Inches(0.3))
        ctf = cbox.text_frame
        ctf.clear()
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = caption
        set_run(cr, size=10, bold=True, color=GRAY)


def slide_locations(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Project Locations", "Nationwide presence — Jakarta HQ & East Kalimantan operations")
    add_footer(slide)

    locations = [
        ("Head Office", "North Jakarta", "Rukan Grand Orchard Square, Jl. Terusan Klp. Hybrida Blok D22, Sukapura, Cilincing, Jakarta 14140"),
        ("Senyiur Hauling Site", "East Kalimantan", "Active coal hauling contract since 2024 — 30 units operating"),
        ("Muara Pahu Project", "East Kalimantan", "0°01'24.3\"N 116°01'45.2\"E — In preparation, go-live 2nd week July 2026 — 23 units"),
    ]

    y = Inches(1.55)
    for i, (title, region, detail) in enumerate(locations):
        accent = TEAL if i == 0 else ORANGE if i == 1 else TEAL
        card = add_rect(slide, Inches(0.65), y, Inches(12.05), Inches(1.45), LIGHT_GRAY)
        card.line.color.rgb = MID_GRAY
        add_rect(slide, Inches(0.65), y, Inches(0.1), Inches(1.45), accent)

        tbox = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.12), Inches(4), Inches(0.35))
        ttf = tbox.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        set_run(tr, size=15, bold=True, color=DARK)

        rbox = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.45), Inches(11), Inches(0.3))
        rtf = rbox.text_frame
        rtf.clear()
        rp = rtf.paragraphs[0]
        rr = rp.add_run()
        rr.text = region
        set_run(rr, size=11, bold=True, color=ORANGE)

        dbox = slide.shapes.add_textbox(Inches(0.95), y + Inches(0.78), Inches(11.2), Inches(0.55))
        dtf = dbox.text_frame
        dtf.word_wrap = True
        dtf.clear()
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = detail
        set_run(dr, size=11, color=GRAY)
        y += Inches(1.62)


def slide_why_us(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, "Why Clients Choose Us", "Proven performance built on safety, reliability, and scale")
    add_footer(slide)

    pillars = [
        ("Safety & Zero Harm Culture", "Strict HSE protocols, driver competency, and continuous safety briefings on every shift."),
        ("Fleet Uptime & Reliability", "Preventive maintenance and rapid breakdown response keep fleet productive."),
        ("Scalable Hauling Capacity", "15 units at start → 30 at Senyiur → 23 planned at Muara Pahu in July 2026."),
        ("Long-Term Client Partnership", "Consistent tonnage delivery, transparent reporting, and professional execution."),
    ]

    for i, (title, desc) in enumerate(pillars):
        col = i % 2
        row = i // 2
        x = Inches(0.65) + col * (Inches(6.05) + Inches(0.35))
        y = Inches(1.45) + row * (Inches(2.55) + Inches(0.2))

        card = add_rect(slide, x, y, Inches(6.05), Inches(2.55), LIGHT_GRAY)
        card.line.color.rgb = MID_GRAY
        add_rect(slide, x, y, Inches(6.05), Inches(0.08), TEAL if i % 2 == 0 else ORANGE)

        tbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.25), Inches(5.55), Inches(0.45))
        ttf = tbox.text_frame
        ttf.clear()
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        set_run(tr, size=14, bold=True, color=DARK)

        dbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.78), Inches(5.55), Inches(1.55))
        dtf = dbox.text_frame
        dtf.word_wrap = True
        dtf.clear()
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = desc
        set_run(dr, size=12, color=GRAY)


def slide_contact(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    add_rect(slide, Inches(0.75), Inches(2.5), Inches(1.2), Inches(0.07), TEAL)

    title = slide.shapes.add_textbox(Inches(0.75), Inches(2.75), Inches(11), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Ready to Partner on Your Next Hauling Project?"
    set_run(r, size=30, bold=True, color=WHITE)

    contacts = [
        "Head Office: Rukan Grand Orchard Square, Jl. Terusan Klp. Hybrida Blok D22, Sukapura, Cilincing, Jakarta 14140",
        "Phone: +62 21 38865143",
        "Email: info@bataramining.com",
        "Website: bataramining.com",
    ]
    y = Inches(3.85)
    for line in contacts:
        box = slide.shapes.add_textbox(Inches(0.75), y, Inches(11.5), Inches(0.4))
        btf = box.text_frame
        btf.clear()
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = line
        set_run(br, size=14, color=TEAL)
        y += Inches(0.45)

    if img("nobgbtrlogo.png").exists():
        slide.shapes.add_picture(str(img("nobgbtrlogo.png")), Inches(10.5), Inches(0.55), width=Inches(1.5))


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_section(prs, "Company Overview", "Coal hauling contractor since 2024")
    slide_overview(prs)
    slide_client(prs)
    slide_section(prs, "Senyiur Project", "Active coal hauling operations")
    slide_senyiur(prs)
    slide_fleet(prs)
    slide_photos_senyiur(prs)
    slide_section(prs, "Our Expertise", "Dedicated coal hauling services")
    slide_hauling(prs)
    slide_section(prs, "Muara Pahu Project", "Upcoming expansion — July 2026")
    slide_muara_pahu(prs)
    slide_photos_muara(prs)
    slide_locations(prs)
    slide_why_us(prs)
    slide_contact(prs)

    prs.save(str(OUT_FILE))
    print(f"Created: {OUT_FILE}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
